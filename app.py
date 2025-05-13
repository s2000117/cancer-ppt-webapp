import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import re

def load_clean_site_table(csv_path: str) -> pd.DataFrame:
    df = pd.read_csv(csv_path, encoding="utf-8-sig", header=0)
    df.columns = df.columns.map(lambda x: re.sub(r'\s+', '', str(x)))
    return df

@st.cache_data
def load_number_data():
    return pd.read_csv("number_utf8.csv", encoding="utf-8-sig")

def highlight_age_groups(ax, ages):
    ranges = {
        "lightcyan": (ages.index("15-19歳"), ages.index("55-59歳")+1),
        "lightblue": (ages.index("60-64歳"), ages.index("70-74歳")+1),
        "steelblue": (ages.index("75-79歳"), ages.index("100歳以上")+1),
    }
    for color, (start, end) in ranges.items():
        ax.axvspan(start - 0.5, end - 0.5, color=color, alpha=0.3)

def create_ppt_slide(prs, fig_path, title_text):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    slide.shapes.add_picture(str(fig_path), Inches(1), Inches(1.5), height=Inches(4.5))

st.title("📊 年齢階級別 がん罹患数 PPT作成 Webアプリ")

site_df = load_clean_site_table("部位別コード表_utf8.csv")
number_df = load_number_data()

years = list(range(2016, 2021))
selected_year = st.selectbox("診断年を選択", years)

sex_options = ["男", "女", "総数"]
selected_sex = st.selectbox("性別を選択", sex_options)

site_list = ["全部位"] + site_df["部位"].tolist()
selected_sites = st.multiselect("部位を選択（複数選択可）", site_list, default=["全部位"])

if st.button("🌐 一括PPTスライド生成"):
    fig_dir = Path("figures")
    fig_dir.mkdir(exist_ok=True)

    age_columns = number_df.columns[number_df.columns.get_loc("0-4歳"):number_df.columns.get_loc("100歳以上") + 1]
    ages = age_columns.tolist()

    prs = Presentation()

    if selected_sites:
        target_sites = site_df["部位"].tolist() if "全部位" in selected_sites else selected_sites
    else:
        st.warning("部位を選択してください")
        target_sites = []

    for site_name in target_sites:
        site_code = site_df[site_df["部位"] == site_name]["コード"].values[0]
        fig, ax = plt.subplots(figsize=(12, 6))
        highlight_age_groups(ax, ages)

        try:
            if selected_sex in ["男", "総数"] and site_name not in ["子宮", "子宮頸部", "子宮体部", "卵巣"]:
                male_row = number_df[(number_df["コード"] == site_code) &
                                     (number_df["性別"] == "男") &
                                     (number_df["診断年"] == selected_year)]
                if not male_row.empty:
                    male_by_age = male_row[age_columns].iloc[0].astype(float)
                    ax.plot(ages, male_by_age, label="男性", color="orange", marker="o")

            if selected_sex in ["女", "総数"] and site_name not in ["前立腺"]:
                female_row = number_df[(number_df["コード"] == site_code) &
                                       (number_df["性別"] == "女") &
                                       (number_df["診断年"] == selected_year)]
                if not female_row.empty:
                    female_by_age = female_row[age_columns].iloc[0].astype(float)
                    ax.plot(ages, female_by_age, label="女性", color="yellow", marker="o")

            ax.set_title(f"{site_name}（{selected_year}年・{selected_sex}）", fontsize=14)
            ax.set_xlabel("年齢階級")
            ax.set_ylabel("罹患数")
            ax.set_xticks(range(len(ages)))
            ax.set_xticklabels(ages, rotation=45)
            ax.legend(loc="upper center", ncol=2)
            fig.tight_layout()

            fig_path = fig_dir / f"{site_code}_{selected_year}_{selected_sex}.png"
            fig.savefig(fig_path)
            plt.close(fig)

            create_ppt_slide(prs, fig_path, f"{site_name}（{selected_year}年・{selected_sex}）")

        except Exception as e:
            st.warning(f"⚠️ {site_name} の処理でエラー: {e}")

    ppt_output_path = f"全部位_{selected_year}_{selected_sex}.pptx"
    prs.save(ppt_output_path)
    with open(ppt_output_path, "rb") as f:
        st.download_button("📥 PPTXをダウンロード", f, file_name=ppt_output_path)
